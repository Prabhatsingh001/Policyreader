import os
import itertools
import asyncio
import json
from typing import List, Optional, Union
from fastapi import FastAPI, Header, HTTPException, status, Depends
from pydantic import BaseModel, Field
from google import genai
from google.genai import types
import httpx
from datetime import datetime
from io import BytesIO
from PyPDF2 import PdfReader, PdfWriter
from docx import Document # Added for DOCX processing
from pptx import Presentation # Added for PPTX processing
from pptx.enum.shapes import MSO_SHAPE_TYPE # Added for PPTX text extraction

# -------- Configuration --------
API_KEYS = []
MAX_CONCURRENT_CALLS = len(API_KEYS)
CLIENTS = [genai.Client(api_key=key) for key in API_KEYS]
TEAM_TOKEN = ""
TIMEOUT_SECONDS = 55
FALLBACK_ANSWER = "I'm sorry, I can't answer that question."
MAX_SPLIT_SIZE_MB = 5 # This applies to PDF splitting. For text/images, we send the whole extracted content.

# A separate API key for the non-document call
NON_DOCUMENT_API_KEY = ""
NON_DOCUMENT_CLIENT = genai.Client(api_key=NON_DOCUMENT_API_KEY)

# HTTP client singleton
HTTP_CLIENT: httpx.AsyncClient

# -------- Models --------
class AnswerWithConfidence(BaseModel):
    answer: str
    confidence: int = Field(..., ge=0, le=100)
    source_type: str = Field(..., description="One of 'document_specific', 'general_knowledge', or 'prohibited'")

class HackRxRequest(BaseModel):
    documents: str # URL to the document
    questions: List[str]

class HackRxResponse(BaseModel):
    answers: List[str]

# Internal model for Gemini's response, including confidence
class InternalGeminiResponse(BaseModel):
    answers: List[AnswerWithConfidence]

# -------- App Init --------
app = FastAPI(root_path="/api/v1")

# -------- Lifecycle Events --------
@app.on_event("startup")
async def startup_event():
    global HTTP_CLIENT
    HTTP_CLIENT = httpx.AsyncClient()

@app.on_event("shutdown")
async def shutdown_event():
    await HTTP_CLIENT.aclose()

# -------- Auth Dependency --------
def verify_token(authorization: str = Header(...)):
    try:
        scheme, token = authorization.split()
    except ValueError:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid auth header")
    if scheme.lower() != "bearer" or token != TEAM_TOKEN:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid token")
    return token

# -------- PDF Splitting Helper --------
async def _split_pdf_by_size(pdf_bytes: bytes) -> List[bytes]:
    """Splits a PDF into smaller chunks of approximately MAX_SPLIT_SIZE_MB."""
    pdf_reader = PdfReader(BytesIO(pdf_bytes))
    total_pages = len(pdf_reader.pages)

    # Calculate average page size to estimate how many pages to put in a chunk
    # Ensure this calculation prevents division by zero if pdf_bytes is empty or pages are 0
    estimated_page_size = len(pdf_bytes) / total_pages if total_pages > 0 else 1
    pages_per_chunk = int((MAX_SPLIT_SIZE_MB * 1024 * 1024) / estimated_page_size)
    if pages_per_chunk == 0: # Handle very small PDFs
        pages_per_chunk = 1

    split_pdfs = []
    for start_page in range(0, total_pages, pages_per_chunk):
        pdf_writer = PdfWriter()
        end_page = min(start_page + pages_per_chunk, total_pages)
        for page_num in range(start_page, end_page):
            pdf_writer.add_page(pdf_reader.pages[page_num])

        output_stream = BytesIO()
        pdf_writer.write(output_stream)
        split_pdfs.append(output_stream.getvalue())

    print(f"📄 PDF split into {len(split_pdfs)} parts.")
    return split_pdfs

# -------- Document Content Extraction Helpers --------
async def _extract_text_from_docx(docx_bytes: bytes) -> str:
    """Extracts text from a DOCX file."""
    try:
        document = Document(BytesIO(docx_bytes))
        full_text = []
        for para in document.paragraphs:
            full_text.append(para.text)
        print("✅ Text extracted from DOCX.")
        return "\n".join(full_text)
    except Exception as e:
        print(f"❌ Failed to extract text from DOCX: {e}")
        return ""

async def _extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """Extracts text from a PPTX file."""
    try:
        presentation = Presentation(BytesIO(pptx_bytes))
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            full_text.append(run.text)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: # Handle grouped shapes
                    for s in shape.shapes:
                        if s.has_text_frame:
                            for paragraph in s.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    full_text.append(run.text)
        print("✅ Text extracted from PPTX.")
        return "\n".join(full_text)
    except Exception as e:
        print(f"❌ Failed to extract text from PPTX: {e}")
        return ""

# -------- Helper to call Gemini (with document or text) --------
async def _call_gemini_with_content(client: genai.Client, content_data: Union[bytes, str], mime_type: str, questions: List[str]) -> Optional[List[AnswerWithConfidence]]:
    """
    Calls Gemini with either PDF bytes, extracted text content, or image bytes.
    content_data: bytes for PDF/Images, str for DOCX/PPTX extracted text.
    mime_type: "application/pdf", "text/plain", "image/jpeg", "image/png".
    """
    instruction = (
        "You are an expert insurance policy assistant. Your primary function is to analyze the provided document content and answer a list of user questions. "
        "For each question, you must provide a JSON object with three fields: 'answer', 'confidence' (0-100), and 'source_type' ('document_specific', 'general_knowledge', or 'prohibited'). "
        "Use 'document_specific' if the answer is in the document, with high confidence. "
        "Use 'general_knowledge' if the question is not in the document but can be answered from your knowledge; confidence must be 0. "
        "Use 'prohibited' for illegal/unethical questions; the answer must be 'I cannot assist with this request as it involves illegal or prohibited actions.' and confidence must be 0."
    )

    if mime_type in ["application/pdf", "image/jpeg", "image/png"]:
        contents = [
            types.Part.from_bytes(data=content_data, mime_type=mime_type),
            instruction,
            f"Questions: {json.dumps(questions)}"
        ]
    elif mime_type == "text/plain":
        contents = [
            types.Part.from_text(text=content_data), # Use from_text for string content
            instruction,
            f"Questions: {json.dumps(questions)}"
        ]
    else:
        print(f"❌ Unsupported MIME type for Gemini content: {mime_type}")
        return None

    def sync_generate():
        return client.models.generate_content(
            model="gemini-2.5-flash",
            contents=contents,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=InternalGeminiResponse,
                thinking_config=types.ThinkingConfig(thinking_budget=0),
            )
        )

    try:
        # Removed inner timeout
        response = asyncio.get_running_loop().run_in_executor(None, sync_generate)
    except Exception as e:
        print(f"❌ Gemini call with content failed: {e}")
        return None

    try:
        parsed: InternalGeminiResponse = (await response).parsed # Await the response here
        # Ensure the number of answers matches the number of questions
        if not isinstance(parsed, InternalGeminiResponse) or len(parsed.answers) != len(questions):
            # If the model returns fewer answers than questions, pad with FALLBACK_ANSWER
            if isinstance(parsed, InternalGeminiResponse) and len(parsed.answers) < len(questions):
                print(f"⚠️ Gemini returned {len(parsed.answers)} answers for {len(questions)} questions. Padding with fallback.")
                padded_answers = parsed.answers + [AnswerWithConfidence(answer=FALLBACK_ANSWER, confidence=0, source_type="general_knowledge") for _ in range(len(questions) - len(parsed.answers))]
                parsed.answers = padded_answers
            else:
                raise ValueError("Invalid or mismatched response schema")
        return parsed.answers
    except Exception as e:
        print(f"❌ Failed to parse Gemini response with content: {e}")
        return None

# -------- Helper to call Gemini (without document) --------
async def _call_gemini_without_document(client: genai.Client, questions: List[str]) -> Optional[List[str]]:
    """A direct Gemini call with just the questions."""
    instruction = (
        "You are an expert assistant. Answer each question in order, under 50 words, using your general knowledge. "
        "Return ONLY a JSON object: {\"answers\": [\"answer1\", \"answer2\", ...]}."
    )

    contents = [
        instruction,
        f"Questions: {json.dumps(questions)}"
    ]

    def sync_generate():
        return client.models.generate_content(
            model="gemini-2.5-flash",
            contents=contents,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=HackRxResponse,
                thinking_config=types.ThinkingConfig(thinking_budget=0),
            )
        )

    try:
        # Removed inner timeout
        response = asyncio.get_running_loop().run_in_executor(None, sync_generate)
        parsed: HackRxResponse = (await response).parsed # Await the response here
        # Ensure the number of answers matches the number of questions
        if not isinstance(parsed, HackRxResponse) or len(parsed.answers) != len(questions):
            if isinstance(parsed, HackRxResponse) and len(parsed.answers) < len(questions):
                print(f"⚠️ Non-document Gemini returned {len(parsed.answers)} answers for {len(questions)} questions. Padding with fallback.")
                padded_answers = parsed.answers + [FALLBACK_ANSWER for _ in range(len(questions) - len(parsed.answers))]
                parsed.answers = padded_answers
            else:
                raise ValueError("Invalid or mismatched response schema for non-document call")
        return parsed.answers
    except Exception as e:
        print(f"❌ Gemini call without document failed: {e}")
        return None

# -------- Endpoint --------
@app.post("/hackrx/run", response_model=HackRxResponse)
async def run_submission(
    payload: HackRxRequest,
    token: str = Depends(verify_token)
):
    print(f"🚀 Processing request for URL: {payload.documents}")

    # Start the non-document call immediately, as it can serve as a fallback
    non_document_answers_task = asyncio.create_task(_call_gemini_without_document(NON_DOCUMENT_CLIENT, payload.questions))

    final_response_from_document_processing: Optional[HackRxResponse] = None
    
    try:
        # Define an inner asynchronous function to encapsulate the document processing pipeline.
        # This allows us to apply a single timeout to the entire process.
        async def _full_document_pipeline():
            doc_bytes = None
            content_type = None
            
            try:
                # Attempt to download the document without an inner timeout
                r = await HTTP_CLIENT.get(payload.documents)
                r.raise_for_status() # Raise an exception for bad status codes (4xx or 5xx)
                doc_bytes = r.content
                content_type = r.headers.get("Content-Type", "").lower()
                
                # Fallback to extension if Content-Type is generic or missing
                if not content_type or "octet-stream" in content_type:
                    if payload.documents.lower().endswith(".pdf"):
                        content_type = "application/pdf"
                    elif payload.documents.lower().endswith(".docx"):
                        content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    elif payload.documents.lower().endswith(".pptx"):
                        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    elif payload.documents.lower().endswith((".jpg", ".jpeg")):
                        content_type = "image/jpeg"
                    elif payload.documents.lower().endswith(".png"):
                        content_type = "image/png"
                
                print(f"✅ Document downloaded successfully. Content-Type: {content_type}")

            except Exception as e:
                print(f"❌ Failed to download document: {e}")
                # If download fails, return None to signal that document processing couldn't proceed
                return None 

            # Determine processing path based on content_type
            if "application/pdf" in content_type:
                return await _process_pdf_document(payload, doc_bytes, non_document_answers_task)
            elif "wordprocessingml.document" in content_type or payload.documents.lower().endswith(".docx"):
                return await _process_docx_document(payload, doc_bytes, non_document_answers_task)
            elif "presentationml.presentation" in content_type or payload.documents.lower().endswith(".pptx"):
                return await _process_pptx_document(payload, doc_bytes, non_document_answers_task)
            elif "image/jpeg" in content_type or "image/png" in content_type:
                return await _process_image_document(payload, doc_bytes, content_type, non_document_answers_task)
            else:
                # If document type is unsupported, explicitly return None to signal fallback to non-doc answers
                print(f"❌ Unsupported document type: {content_type}. Will only use general knowledge answers.")
                return None # Signal that document processing should be skipped

        # Attempt the full document processing pipeline with the overall TIMEOUT_SECONDS
        final_response_from_document_processing = await asyncio.wait_for(_full_document_pipeline(), timeout=TIMEOUT_SECONDS)
        
        # If document processing was successful and returned a response, return it directly
        if final_response_from_document_processing is not None:
            print("✅ Document-specific processing completed successfully.")
            return final_response_from_document_processing

    except asyncio.TimeoutError:
        print(f"❌ Overall request timed out after {TIMEOUT_SECONDS} seconds.")
    except Exception as e:
        print(f"❌ An unexpected error occurred during document processing: {e}")

    # -------- Fallback Logic --------
    # This section is reached if:
    # 1. The document processing timed out (asyncio.TimeoutError).
    # 2. The document processing encountered an unhandled exception.
    # 3. The document download failed (final_response_from_document_processing was None).
    # 4. The document type was unsupported (final_response_from_document_processing was None).

    print("Attempting to return fallback answers from non-document call or generic fallback.")
    
    # Try to get answers from the non-document task
    if not non_document_answers_task.done():
        try:
            # Removed inner timeout for non_document_answers_task
            fallback_answers = await non_document_answers_task
            print("✅ Returning answers from non-document call (fallback).")
            return HackRxResponse(answers=fallback_answers)
        except Exception as e:
            print(f"❌ Non-document call also failed during fallback: {e}.")
            # Fall through to generic fallback if non-document task also fails
    elif non_document_answers_task.done():
        try:
            # If the non-document task already completed, get its result
            fallback_answers = non_document_answers_task.result()
            print("✅ Returning answers from non-document call (already completed).")
            return HackRxResponse(answers=fallback_answers)
        except Exception as e:
            print(f"❌ Non-document call result retrieval failed: {e}.")
            # Fall through to generic fallback if non-document task result retrieval fails

    # Ultimate fallback: return generic answers if all else fails
    answers = [FALLBACK_ANSWER for _ in payload.questions]
    print("❌ All attempts failed. Returning generic fallback answers.")
    return HackRxResponse(answers=answers)


async def _process_pdf_document(payload: HackRxRequest, pdf_bytes: bytes, non_document_answers_task: asyncio.Task) -> HackRxResponse:
    """Processes PDF documents."""
    split_pdfs = await _split_pdf_by_size(pdf_bytes)
    if not split_pdfs: # Handle case where PDF splitting results in no parts (e.g., empty PDF)
        print("❌ PDF splitting resulted in no parts. Falling back.")
        # Return a response indicating failure, which will be caught by _full_document_pipeline
        return HackRxResponse(answers=[FALLBACK_ANSWER for _ in payload.questions])

    # Truncate splits if they exceed the number of available API keys
    if len(split_pdfs) > MAX_CONCURRENT_CALLS:
        split_pdfs = split_pdfs[:MAX_CONCURRENT_CALLS]
        print(f"⚠️ Truncated PDF splits to match the number of API keys ({MAX_CONCURRENT_CALLS}).")

    print(f"✨ Making {len(split_pdfs)} concurrent API calls for PDF.")
    tasks = []
    for i, pdf_part in enumerate(split_pdfs):
        client = CLIENTS[i % len(CLIENTS)] # Use modulo to cycle through clients if more splits than clients
        tasks.append(_call_gemini_with_content(client, pdf_part, "application/pdf", payload.questions))

    return await _collate_and_respond(payload, tasks, non_document_answers_task)

async def _process_docx_document(payload: HackRxRequest, docx_bytes: bytes, non_document_answers_task: asyncio.Task) -> HackRxResponse:
    """Processes DOCX documents by extracting text."""
    extracted_text = await _extract_text_from_docx(docx_bytes)
    if not extracted_text.strip(): # Check if extracted text is empty or just whitespace
        print("❌ No meaningful text extracted from DOCX. Falling back.")
        return HackRxResponse(answers=[FALLBACK_ANSWER for _ in payload.questions])

    # For text, we send the whole extracted text as one part.
    # If the text is extremely large, this might hit token limits.
    client = CLIENTS[0] # Use the first client for text processing
    tasks = [_call_gemini_with_content(client, extracted_text, "text/plain", payload.questions)]
    
    print(f"✨ Making 1 API call for DOCX text extraction.")
    return await _collate_and_respond(payload, tasks, non_document_answers_task)

async def _process_pptx_document(payload: HackRxRequest, pptx_bytes: bytes, non_document_answers_task: asyncio.Task) -> HackRxResponse:
    """Processes PPTX documents by extracting text."""
    extracted_text = await _extract_text_from_pptx(pptx_bytes)
    if not extracted_text.strip(): # Check if extracted text is empty or just whitespace
        print("❌ No meaningful text extracted from PPTX. Falling back.")
        return HackRxResponse(answers=[FALLBACK_ANSWER for _ in payload.questions])

    client = CLIENTS[0] # Use the first client for text processing
    tasks = [_call_gemini_with_content(client, extracted_text, "text/plain", payload.questions)]
    
    print(f"✨ Making 1 API call for PPTX text extraction.")
    return await _collate_and_respond(payload, tasks, non_document_answers_task)

async def _process_image_document(payload: HackRxRequest, image_bytes: bytes, image_mime_type: str, non_document_answers_task: asyncio.Task) -> HackRxResponse:
    """Processes image documents (JPG, PNG)."""
    if not image_bytes:
        print("❌ No image bytes provided. Falling back.")
        return HackRxResponse(answers=[FALLBACK_ANSWER for _ in payload.questions])

    # For images, we send the whole image as one part.
    client = CLIENTS[0] # Use the first client for image processing
    tasks = [_call_gemini_with_content(client, image_bytes, image_mime_type, payload.questions)]
    
    print(f"✨ Making 1 API call for {image_mime_type} image processing.")
    return await _collate_and_respond(payload, tasks, non_document_answers_task)


async def _collate_and_respond(payload: HackRxRequest, tasks: List[asyncio.Task], non_document_answers_task: asyncio.Task) -> HackRxResponse:
    """Collates responses from Gemini calls and forms the final response."""
    all_responses = await asyncio.gather(*tasks, return_exceptions=True)
    
    # Await the non-document call if it hasn't finished already
    if not non_document_answers_task.done():
        non_document_answers = await non_document_answers_task
    else:
        non_document_answers = non_document_answers_task.result()
    
    print("✅ All API calls completed (or failed).")

    final_answers = []
    num_questions = len(payload.questions)
    
    for i in range(num_questions):
        prohibited_keywords = ["fraudulent", "illegal", "forged", "manipulate"]
        q_lower = payload.questions[i].lower()
        if any(word in q_lower for word in prohibited_keywords):
            final_answers.append("I cannot assist with this request as it involves illegal or prohibited actions.")
            continue

        best_answer = FALLBACK_ANSWER
        highest_confidence = -1
        doc_specific_count = 0
        general_knowledge_count = 0
        
        for response in all_responses:
            if isinstance(response, Exception) or response is None:
                continue
            
            # Ensure the response list has enough elements for the current question index
            if i < len(response):
                answer_item = response[i]
                
                if answer_item.source_type == "document_specific":
                    doc_specific_count += 1
                    if answer_item.confidence > highest_confidence:
                        highest_confidence = answer_item.confidence
                        best_answer = answer_item.answer
                elif answer_item.source_type == "general_knowledge":
                    general_knowledge_count += 1
                    
        # Decision logic
        if doc_specific_count > general_knowledge_count and highest_confidence > 0:
            final_answers.append(best_answer)
        elif non_document_answers and i < len(non_document_answers):
            final_answers.append(non_document_answers[i])
        else:
            final_answers.append(FALLBACK_ANSWER)

    print("✅ Final answers selected based on confidence scores and source type.")

    try:
        with open("queries.log", "a", encoding="utf-8") as log_file:
            log_entry = {
                "timestamp": datetime.utcnow().isoformat(),
                "document_url": payload.documents, # Changed from pdf_url
                "questions": payload.questions,
                "answers": final_answers
            }
            log_file.write(json.dumps(log_entry, ensure_ascii=False) + "\n")
        print("📝 Log entry created.")
    except Exception as e:
        print(f"❌ Failed to write to log file: {e}")

    return HackRxResponse(answers=final_answers)

# -------- Uvicorn Entry Point --------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
